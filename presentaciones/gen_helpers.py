"""Helpers visuales VotoAfin — v2 (visual-first, minimo texto en slides)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR

# -- Brand Colors -------------------------------------------------------------
HERO    = RGBColor(0x1C, 0x3A, 0x52)
PRIMARY = RGBColor(0x2E, 0x5F, 0x7E)
PRILIT  = RGBColor(0x4A, 0x9B, 0xBF)
SEC     = RGBColor(0x7B, 0xA0, 0x98)
ACCENT  = RGBColor(0x3A, 0x9E, 0x7A)
SURFACE = RGBColor(0xF7, 0xF8, 0xF7)
DARK2   = RGBColor(0x24, 0x4A, 0x63)
TEXT    = RGBColor(0x1F, 0x2A, 0x35)
TEXT2   = RGBColor(0x4A, 0x55, 0x68)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DANGER  = RGBColor(0xDC, 0x26, 0x26)
AMBER   = RGBColor(0xD9, 0x77, 0x06)
BORDER  = RGBColor(0xCB, 0xD1, 0xCD)

# Eje colors
EJE_ECO   = RGBColor(0x15, 0x80, 0x3D)
EJE_SOC   = RGBColor(0x1D, 0x4E, 0xD8)
EJE_AMB   = RGBColor(0x06, 0x5F, 0x46)
EJE_SEG   = RGBColor(0x92, 0x40, 0x0E)
EJE_DDHH  = RGBColor(0x7C, 0x3A, 0xED)
EJE_INT   = RGBColor(0x0E, 0x75, 0x90)
EJE_INST  = RGBColor(0x92, 0x60, 0x0E)

# -- Slide dimensions 16:9 ----------------------------------------------------
W  = Inches(13.333)
H  = Inches(7.5)
MG = Inches(0.7)
SW = W - 2 * MG   # safe width


def new_prs():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)


def _rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    return s


def _oval(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(9, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def _p1(tf, text, size, bold, color, align=PP_ALIGN.LEFT, italic=False, space=0):
    p = tf.add_paragraph()
    p.text = text
    p.alignment = align
    if space:
        p.space_before = Pt(space)
    r = p.runs[0]
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return p


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _top_bar(slide, color=PRIMARY):
    _rect(slide, 0, 0, W, Inches(0.07), color)


def _accent_line(slide, x=MG, y=Inches(1.25), w=Inches(1.5)):
    _rect(slide, x, y, w, Inches(0.04), ACCENT)


# =============================================================================
# SLIDE TYPES
# =============================================================================

def dark_cover(prs, title, subtitle, label="", notes=""):
    """Portada oscura hero."""
    sl = _blank(prs)
    _bg(sl, HERO)
    _rect(sl, 0, 0, W, Inches(0.05), ACCENT)

    if label:
        tb = _tb(sl, MG, Inches(1.5), SW, Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]; p.text = label.upper()
        r = p.runs[0]; r.font.size = Pt(10); r.font.bold = True
        r.font.color.rgb = ACCENT

    _rect(sl, MG, Inches(2.05), Inches(0.05), Inches(2.8), ACCENT)

    tb2 = _tb(sl, MG + Inches(0.2), Inches(2.1), SW - Inches(0.2), Inches(2.6))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.text = title
    r2 = p2.runs[0]; r2.font.size = Pt(42); r2.font.bold = True
    r2.font.color.rgb = WHITE

    tb3 = _tb(sl, MG + Inches(0.2), Inches(4.9), SW - Inches(0.2), Inches(1.2))
    tf3 = tb3.text_frame; tf3.word_wrap = True
    p3 = tf3.paragraphs[0]; p3.text = subtitle
    r3 = p3.runs[0]; r3.font.size = Pt(17); r3.font.bold = False
    r3.font.color.rgb = PRILIT

    tb4 = _tb(sl, MG, Inches(6.85), SW, Inches(0.35))
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]; p4.alignment = PP_ALIGN.RIGHT
    p4.text = "VotoAfin  --  UTFSM  --  Jenifer Castillo  --  2026"
    r4 = p4.runs[0]; r4.font.size = Pt(9); r4.font.color.rgb = SEC

    if notes:
        _notes(sl, notes)


def section_break(prs, number, title, subtitle="", notes=""):
    """Quiebre de seccion visual."""
    sl = _blank(prs)
    _bg(sl, PRIMARY)
    _rect(sl, 0, 0, W, Inches(0.05), ACCENT)

    # Big number ghost
    tb = _tb(sl, Inches(9.5), Inches(0.5), Inches(3.5), Inches(5))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = number; p.alignment = PP_ALIGN.RIGHT
    r = p.runs[0]; r.font.size = Pt(160); r.font.bold = True
    r.font.color.rgb = DARK2

    _rect(sl, MG, Inches(2.8), Inches(0.06), Inches(1.8), ACCENT)

    tb2 = _tb(sl, MG + Inches(0.2), Inches(2.8), Inches(8.5), Inches(1.8))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.text = title
    r2 = p2.runs[0]; r2.font.size = Pt(38); r2.font.bold = True
    r2.font.color.rgb = WHITE

    if subtitle:
        tb3 = _tb(sl, MG + Inches(0.2), Inches(4.7), Inches(8.5), Inches(0.8))
        tf3 = tb3.text_frame; tf3.word_wrap = True
        p3 = tf3.paragraphs[0]; p3.text = subtitle
        r3 = p3.runs[0]; r3.font.size = Pt(16); r3.font.color.rgb = PRILIT

    if notes:
        _notes(sl, notes)


def stat_slide(prs, title, stats, notes=""):
    """4 estadisticas grandes. stats = list of (number, label, color)."""
    sl = _blank(prs)
    _bg(sl, HERO)
    _rect(sl, 0, 0, W, Inches(0.06), ACCENT)

    tb = _tb(sl, MG, Inches(0.25), SW, Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = title
    r = p.runs[0]; r.font.size = Pt(20); r.font.bold = True
    r.font.color.rgb = PRILIT

    n = len(stats)
    cw = SW / n
    for i, (num, lbl, col) in enumerate(stats):
        x = MG + i * cw
        # card bg
        _rect(sl, x + Inches(0.1), Inches(1.15), cw - Inches(0.2), Inches(5.8),
              DARK2)
        # number
        tbn = _tb(sl, x + Inches(0.1), Inches(1.5), cw - Inches(0.2), Inches(2.8))
        tfn = tbn.text_frame
        pn = tfn.paragraphs[0]; pn.text = num; pn.alignment = PP_ALIGN.CENTER
        rn = pn.runs[0]; rn.font.size = Pt(64); rn.font.bold = True
        rn.font.color.rgb = col
        # label
        tbl = _tb(sl, x + Inches(0.15), Inches(4.3), cw - Inches(0.3), Inches(2.2))
        tfl = tbl.text_frame; tfl.word_wrap = True
        pl = tfl.paragraphs[0]; pl.text = lbl; pl.alignment = PP_ALIGN.CENTER
        rl = pl.runs[0]; rl.font.size = Pt(13); rl.font.bold = False
        rl.font.color.rgb = PRILIT

    if notes:
        _notes(sl, notes)


def quote_slide(prs, quote, author, notes=""):
    """Cita grande sobre fondo oscuro."""
    sl = _blank(prs)
    _bg(sl, DARK2)
    _rect(sl, 0, 0, W, Inches(0.05), ACCENT)

    # quote mark
    tq = _tb(sl, MG, Inches(0.3), Inches(1.2), Inches(1.2))
    tfq = tq.text_frame
    pq = tfq.paragraphs[0]; pq.text = "\u201c"
    rq = pq.runs[0]; rq.font.size = Pt(80); rq.font.bold = True
    rq.font.color.rgb = ACCENT

    tb = _tb(sl, MG, Inches(1.35), SW, Inches(4.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = quote; p.alignment = PP_ALIGN.LEFT
    r = p.runs[0]; r.font.size = Pt(22); r.font.italic = True
    r.font.color.rgb = WHITE

    _rect(sl, MG, Inches(5.75), Inches(2.0), Inches(0.04), ACCENT)

    tba = _tb(sl, MG, Inches(5.9), SW, Inches(0.55))
    tfa = tba.text_frame
    pa = tfa.paragraphs[0]; pa.text = f"-- {author}"; pa.alignment = PP_ALIGN.LEFT
    ra = pa.runs[0]; ra.font.size = Pt(12); ra.font.bold = True
    ra.font.color.rgb = SEC

    if notes:
        _notes(sl, notes)


def tight_bullets(prs, title, bullets, notes="", dark=False):
    """Max 5 bullets cortos (<=8 palabras c/u). NO parrafos."""
    sl = _blank(prs)
    bg_col = HERO if dark else SURFACE
    _bg(sl, bg_col)
    _rect(sl, 0, 0, W, Inches(0.07), PRIMARY if not dark else ACCENT)

    t_col  = WHITE if dark else HERO
    b_col  = PRILIT if dark else TEXT

    tb = _tb(sl, MG, Inches(0.3), SW, Inches(0.82))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    r = p.runs[0]; r.font.size = Pt(28); r.font.bold = True
    r.font.color.rgb = t_col

    _rect(sl, MG, Inches(1.22), Inches(1.2), Inches(0.04),
          ACCENT if not dark else WHITE)

    start_y = Inches(1.45)
    gap = (H - start_y - Inches(0.4)) / max(len(bullets), 1)
    gap = min(gap, Inches(1.1))

    for i, bullet in enumerate(bullets):
        y = start_y + i * gap
        # dot
        _oval(sl, MG, y + Inches(0.15), Inches(0.18), Inches(0.18),
              ACCENT if not dark else PRILIT)
        # text
        tbb = _tb(sl, MG + Inches(0.35), y, SW - Inches(0.35), gap)
        tfb = tbb.text_frame; tfb.word_wrap = True
        pb = tfb.paragraphs[0]; pb.text = bullet
        rb = pb.runs[0]; rb.font.size = Pt(19); rb.font.bold = False
        rb.font.color.rgb = b_col

    if notes:
        _notes(sl, notes)


def single_message(prs, message, sub="", dark=True, notes=""):
    """Un mensaje grande. Sin bullets. Visual de impacto."""
    sl = _blank(prs)
    _bg(sl, HERO if dark else SURFACE)
    _rect(sl, 0, 0, W, Inches(0.05), ACCENT)

    msg_y = Inches(2.1) if sub else Inches(2.5)
    tb = _tb(sl, MG, msg_y, SW, Inches(2.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = message; p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]; r.font.size = Pt(36); r.font.bold = True
    r.font.color.rgb = WHITE if dark else HERO

    if sub:
        tb2 = _tb(sl, MG, Inches(4.7), SW, Inches(1.4))
        tf2 = tb2.text_frame; tf2.word_wrap = True
        p2 = tf2.paragraphs[0]; p2.text = sub; p2.alignment = PP_ALIGN.CENTER
        r2 = p2.runs[0]; r2.font.size = Pt(18); r2.font.bold = False
        r2.font.color.rgb = PRILIT if dark else TEXT2

    if notes:
        _notes(sl, notes)


def big_question(prs, question, answer="", notes=""):
    """Pregunta grande centrada sobre fondo oscuro."""
    sl = _blank(prs)
    _bg(sl, HERO)
    _rect(sl, 0, 0, W, Inches(0.05), ACCENT)

    # decorative line
    _rect(sl, Inches(3.5), Inches(1.8), Inches(6.3), Inches(0.04), PRILIT)

    tb = _tb(sl, MG, Inches(2.0), SW, Inches(2.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = question; p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]; r.font.size = Pt(30); r.font.bold = True
    r.font.color.rgb = WHITE

    if answer:
        _rect(sl, Inches(3.5), Inches(4.9), Inches(6.3), Inches(0.04), ACCENT)
        tb2 = _tb(sl, MG, Inches(5.1), SW, Inches(1.5))
        tf2 = tb2.text_frame; tf2.word_wrap = True
        p2 = tf2.paragraphs[0]; p2.text = answer; p2.alignment = PP_ALIGN.CENTER
        r2 = p2.runs[0]; r2.font.size = Pt(22); r2.font.bold = False
        r2.font.color.rgb = ACCENT

    if notes:
        _notes(sl, notes)


def timeline_slide(prs, title, items, notes=""):
    """Timeline horizontal. items = list of (year, label, color).
    Max 5 items recomendado."""
    sl = _blank(prs)
    _bg(sl, SURFACE)
    _rect(sl, 0, 0, W, Inches(0.07), PRIMARY)

    tb = _tb(sl, MG, Inches(0.28), SW, Inches(0.75))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = title
    r = p.runs[0]; r.font.size = Pt(26); r.font.bold = True
    r.font.color.rgb = HERO

    n = len(items)
    line_y = Inches(3.5)
    seg_w  = SW / n
    r_w    = Inches(0.34)

    # horizontal line
    _rect(sl, MG, line_y - Inches(0.03), SW, Inches(0.06), BORDER)

    for i, (year, label, col) in enumerate(items):
        cx = MG + i * seg_w + seg_w / 2

        # circle
        _oval(sl, cx - r_w / 2, line_y - r_w / 2, r_w, r_w, col)

        # year above
        tby = _tb(sl, cx - Inches(0.7), line_y - Inches(0.9), Inches(1.4), Inches(0.5))
        tfy = tby.text_frame
        py = tfy.paragraphs[0]; py.text = year; py.alignment = PP_ALIGN.CENTER
        ry = py.runs[0]; ry.font.size = Pt(14); ry.font.bold = True
        ry.font.color.rgb = col

        # label below (up to 2 lines)
        tbl = _tb(sl, cx - Inches(0.95), line_y + Inches(0.4),
                  Inches(1.9), Inches(2.4))
        tfl = tbl.text_frame; tfl.word_wrap = True
        pl = tfl.paragraphs[0]; pl.text = label; pl.alignment = PP_ALIGN.CENTER
        rl = pl.runs[0]; rl.font.size = Pt(13); rl.font.bold = False
        rl.font.color.rgb = TEXT2

    if notes:
        _notes(sl, notes)


def three_cards(prs, title, cards, notes=""):
    """3 tarjetas de contenido. cards = list of (header_text, body, color)."""
    sl = _blank(prs)
    _bg(sl, SURFACE)
    _rect(sl, 0, 0, W, Inches(0.07), PRIMARY)

    tb = _tb(sl, MG, Inches(0.28), SW, Inches(0.75))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = title
    r = p.runs[0]; r.font.size = Pt(26); r.font.bold = True
    r.font.color.rgb = HERO

    card_w = (SW - Inches(0.4)) / 3
    card_h = Inches(5.4)
    top_h  = Inches(1.3)

    for i, (hdr, body, col) in enumerate(cards):
        x = MG + i * (card_w + Inches(0.2))
        y = Inches(1.3)

        # colored top
        _rect(sl, x, y, card_w, top_h, col)
        tbh = _tb(sl, x + Inches(0.12), y + Inches(0.15),
                  card_w - Inches(0.24), top_h - Inches(0.15))
        tfh = tbh.text_frame; tfh.word_wrap = True
        ph = tfh.paragraphs[0]; ph.text = hdr; ph.alignment = PP_ALIGN.CENTER
        rh = ph.runs[0]; rh.font.size = Pt(16); rh.font.bold = True
        rh.font.color.rgb = WHITE

        # white body
        _rect(sl, x, y + top_h, card_w, card_h - top_h, WHITE, BORDER)
        tbb = _tb(sl, x + Inches(0.15), y + top_h + Inches(0.2),
                  card_w - Inches(0.3), card_h - top_h - Inches(0.3))
        tfb = tbb.text_frame; tfb.word_wrap = True
        pb = tfb.paragraphs[0]; pb.text = body; pb.alignment = PP_ALIGN.LEFT
        rb = pb.runs[0]; rb.font.size = Pt(13); rb.font.bold = False
        rb.font.color.rgb = TEXT

    if notes:
        _notes(sl, notes)


def before_after(prs, title, left_title, left_items,
                 right_title, right_items, notes="",
                 left_col=DANGER, right_col=ACCENT):
    """Comparacion visual izquierda / derecha. Max 5 items cada lado."""
    sl = _blank(prs)
    _bg(sl, SURFACE)
    _rect(sl, 0, 0, W, Inches(0.07), PRIMARY)

    tb = _tb(sl, MG, Inches(0.28), SW, Inches(0.75))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = title
    r = p.runs[0]; r.font.size = Pt(26); r.font.bold = True
    r.font.color.rgb = HERO

    col_w  = (SW - Inches(0.25)) / 2
    panels = [
        (MG,                      left_title,  left_items,  left_col),
        (MG + col_w + Inches(0.25), right_title, right_items, right_col),
    ]

    for x, hdr, items, col in panels:
        # header bar
        _rect(sl, x, Inches(1.3), col_w, Inches(0.55), col)
        tbh = _tb(sl, x + Inches(0.15), Inches(1.35),
                  col_w - Inches(0.3), Inches(0.45))
        tfh = tbh.text_frame
        ph = tfh.paragraphs[0]; ph.text = hdr; ph.alignment = PP_ALIGN.CENTER
        rh = ph.runs[0]; rh.font.size = Pt(14); rh.font.bold = True
        rh.font.color.rgb = WHITE

        # body area
        _rect(sl, x, Inches(1.85), col_w, Inches(5.3), WHITE, BORDER)

        gap = Inches(0.85)
        for j, item in enumerate(items[:5]):
            iy = Inches(2.05) + j * gap
            _oval(sl, x + Inches(0.2), iy + Inches(0.12),
                  Inches(0.16), Inches(0.16), col)
            tbi = _tb(sl, x + Inches(0.46), iy,
                      col_w - Inches(0.56), gap - Inches(0.05))
            tfi = tbi.text_frame; tfi.word_wrap = True
            pi = tfi.paragraphs[0]; pi.text = item
            ri = pi.runs[0]; ri.font.size = Pt(14); ri.font.bold = False
            ri.font.color.rgb = TEXT

    if notes:
        _notes(sl, notes)


def flow_steps(prs, title, steps, notes=""):
    """Flujo de pasos. steps = list of (num_str, label, sub). Max 7."""
    sl = _blank(prs)
    _bg(sl, SURFACE)
    _rect(sl, 0, 0, W, Inches(0.07), PRIMARY)

    tb = _tb(sl, MG, Inches(0.28), SW, Inches(0.75))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = title
    r = p.runs[0]; r.font.size = Pt(26); r.font.bold = True
    r.font.color.rgb = HERO

    n = len(steps)
    step_w = SW / n

    for i, (num, label, sub) in enumerate(steps):
        x  = MG + i * step_w
        cx = x + step_w / 2
        circle_r = Inches(0.55)
        cy = Inches(2.8)

        # connector line (not for last)
        if i < n - 1:
            _rect(sl, cx, cy, step_w, Inches(0.04), BORDER)

        col = [ACCENT, PRIMARY, PRILIT, SEC, AMBER, DANGER, DARK2][i % 7]
        _oval(sl, cx - circle_r / 2, cy - circle_r / 2,
              circle_r, circle_r, col)

        # number inside circle
        tbn = _tb(sl, cx - circle_r / 2, cy - circle_r / 2,
                  circle_r, circle_r)
        tfn = tbn.text_frame
        pn = tfn.paragraphs[0]; pn.text = num; pn.alignment = PP_ALIGN.CENTER
        rn = pn.runs[0]; rn.font.size = Pt(18); rn.font.bold = True
        rn.font.color.rgb = WHITE

        # label
        tbl = _tb(sl, x + Inches(0.08), cy + Inches(0.65),
                  step_w - Inches(0.16), Inches(1.1))
        tfl = tbl.text_frame; tfl.word_wrap = True
        pl = tfl.paragraphs[0]; pl.text = label; pl.alignment = PP_ALIGN.CENTER
        rl = pl.runs[0]; rl.font.size = Pt(13); rl.font.bold = True
        rl.font.color.rgb = HERO

        if sub:
            tbs = _tb(sl, x + Inches(0.08), cy + Inches(1.75),
                      step_w - Inches(0.16), Inches(1.2))
            tfs = tbs.text_frame; tfs.word_wrap = True
            ps = tfs.paragraphs[0]; ps.text = sub; ps.alignment = PP_ALIGN.CENTER
            rs = ps.runs[0]; rs.font.size = Pt(11); rs.font.bold = False
            rs.font.color.rgb = TEXT2

    if notes:
        _notes(sl, notes)


def two_col_tight(prs, title, left_title, left_items,
                  right_title, right_items, notes=""):
    """Dos columnas con max 5 bullets cada una."""
    sl = _blank(prs)
    _bg(sl, SURFACE)
    _rect(sl, 0, 0, W, Inches(0.07), PRIMARY)

    tb = _tb(sl, MG, Inches(0.28), SW, Inches(0.75))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = title
    r = p.runs[0]; r.font.size = Pt(26); r.font.bold = True
    r.font.color.rgb = HERO

    col_w = (SW - Inches(0.3)) / 2
    sep_x = MG + col_w + Inches(0.15)
    _rect(sl, sep_x - Inches(0.015), Inches(1.3), Inches(0.03), Inches(5.8), BORDER)

    for i, (hdr, items, x) in enumerate([
        (left_title,  left_items,  MG),
        (right_title, right_items, sep_x + Inches(0.15)),
    ]):
        # header
        tbh = _tb(sl, x, Inches(1.3), col_w, Inches(0.5))
        tfh = tbh.text_frame
        ph = tfh.paragraphs[0]; ph.text = hdr
        rh = ph.runs[0]; rh.font.size = Pt(15); rh.font.bold = True
        rh.font.color.rgb = PRIMARY

        _rect(sl, x, Inches(1.82), Inches(0.9), Inches(0.04), ACCENT)

        gap = Inches(0.88)
        for j, item in enumerate(items[:5]):
            iy = Inches(2.05) + j * gap
            _oval(sl, x, iy + Inches(0.13), Inches(0.16), Inches(0.16), ACCENT)
            tbi = _tb(sl, x + Inches(0.28), iy, col_w - Inches(0.28), gap)
            tfi = tbi.text_frame; tfi.word_wrap = True
            pi = tfi.paragraphs[0]; pi.text = item
            ri = pi.runs[0]; ri.font.size = Pt(14); ri.font.color.rgb = TEXT

    if notes:
        _notes(sl, notes)
